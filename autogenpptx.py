from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import argparse
import sys
import signal

parser = argparse.ArgumentParser(
    description="This tool will help you create" +
    " your presentation slide fast and easy. " +
    "Created by Chanakan Mungtin"
)
parser.add_argument(
    "-i", "--input", help="Specify powerpoint draft file " +
    "created by format within 'format.txt'", nargs='?',
    default=None, required=True
)
parser.add_argument(
    "-p", "--pages", help="Specify how many page will be in " +
    " the presentation -- DEFAULT 1", nargs='?', default=0,
    required=True, type=int
)
parser.add_argument(
    "-o", "--output", help="Specify ouput filename for the " +
    "ouput file. -- DEFAULT output.pptx", nargs='?',
    default="output.pptx", required=False, type=str
)

args = parser.parse_args()

outfilename = str(args.output)


def inputFileRead(filename):
    try:
        f = open(filename, "r")
        fcontent = f.readline()
        return fcontent, f
    except IOError:
        print("Error! File inaccessible")
        sys.exit(1)


def main():
    if args.input is not None:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        content_slide_layout = prs.slide_layouts[1]
        title_slide = prs.slides.add_slide(title_slide_layout)
        content_slide = []
        if args.pages is None:
            print("pages flag should be number.")
            sys.exit(1)
        for i in range(int(args.pages) - 1):
            content_slide.append(prs.slides.add_slide(content_slide_layout))
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        draft, f = inputFileRead(str(args.input))
        createPPTX(title, subtitle, draft, f, prs, content_slide)
    else:
        parser.print_help()


def createPPTX(title, subtitle, draft, f, prs, content_slide):
    signCount = 0
    callCount = 0
    while draft:
        fileContent = str(draft.strip())
        if "#" in fileContent:
            try:
                type(PCDraft)
            except NameError:
                title.text = str(fileContent.replace("#", ""))
                prs.save(outfilename)
            else:
                try:
                    content_title = content_slide[PCDraft].shapes.title
                    content_title.text = str(fileContent.replace("#", ""))
                except(IndexError):
                    print(
                        "Pages is not equal to section in draft file.",
                        "exiting..."
                    )
                    sys.exit(1)
                except(UnboundLocalError):
                    print(
                        "Slide did not have content pages but the draft ",
                        "file specify one. Existing."
                    )
                    sys.exit(1)
                prs.save(outfilename)
        elif "##" in fileContent:
            subtitle.text = str(fileContent.replace("##", ""))
            prs.save(outfilename)
        elif "\newpage" in fileContent:
            callCount += 1
            if args.pages <= 1:
                pass
            else:
                if callCount > 1:
                    PCDraft = signCount + 1
                    signCount += 1
                else:
                    PCDraft = 0

#        elif "content_title>" in fileContent:
#            try:
#                content_title = content_slide[PCDraft].shapes.title
#                content_title.text = str(fileContent.replace("%%%", ""))
#            except(IndexError):
#                print(
#                    "Pages is not equal to section in draft file.",
#                    "exiting..."
#                )
#                sys.exit(1)
#            except(UnboundLocalError):
#                print(
#                    "Slide did not have content pages but the draft ",
#                    "file specify one. Existing."
#                )
#                sys.exit(1)
#            prs.save(outfilename)

        elif "p>" in fileContent:
            try:
                content_con = content_slide[PCDraft].shapes.placeholders[1]
                content_con.text = str(fileContent.replace("p>", ""))
            except(IndexError):
                print("Pages is not enough. exiting...")
                sys.exit(1)
            prs.save(outfilename)
        elif "bgcl>" in fileContent:
            background = content_slide[PCDraft].background
            try:
                R, G, B = str(fileContent.replace("bgcl>", "")).split(",")
            except(ValueError):
                print(
                    "background value should be RGB value seperate ",
                    "by comma. Example: bgcl>255,255,255"
                )
                sys.exit(1)
            fill = background.fill
            fill.solid()
            try:
                fill.fore_color.rgb = RGBColor(int(R), int(G), int(B))
            except(ValueError):
                print(
                    "background value should be RGB value seperate by ",
                    "comma. Example: bgcl>255,255,255"
                )
                sys.exit(1)
            prs.save(outfilename)
        elif "img>" in fileContent:
            con_shapes = content_slide[PCDraft].shapes
            try:
                imgpath, x1, y1, x2, y2 = str(
                    fileContent.replace("img>",
                                        "")
                ).split(",")
            except(ValueError):
                print(
                    "image value should be image path, Position of ",
                    "image horizontal, Position of image verticle, ",
                    "size of image horizontal, size of image verticle in Inches."
                )
                sys.exit(1)
            try:
                con_shapes.add_picture(
                    str(imgpath), Inches(int(x1)),
                    Inches(int(y1)), Inches(int(x2)),
                    Inches(int(y2))
                )
            except(ValueError):
                print(
                    "image value should be image path, Position of ",
                    "image horizontal, Position of image verticle, ",
                    "size of image horizontal, size of image verticle in Inches."
                )
                sys.exit(1)
            except(FileNotFoundError):
                print("image file did not exist. Please try again.")
                sys.exit(1)
            prs.save(outfilename)
        elif "\n" or "\r\n" in fileContent:
            pass
        else:
            print(
                "draft file is corrupted. Check the format.txt file ",
                "and try again."
            )
        draft = f.readline()


if __name__ == "__main__":
    main()
